#!/usr/bin/env python3
"""
File to Markdown Converter
Converts Word (.docx), PDF (.pdf), and PowerPoint (.ppt, .pptx) files to Markdown format.
"""

import os
import sys
import argparse
import tempfile
import uuid
from pathlib import Path


def convert_docx_to_markdown(file_path, image_dir=None):
    """Convert a Word document (.docx) to Markdown.

    If image_dir is provided, embedded images are exported there and referenced
    as Markdown images at the end of the document.
    """
    try:
        from docx import Document
        from docx.oxml.text.paragraph import CT_P
        from docx.oxml.table import CT_Tbl
        from docx.table import Table
        from docx.text.paragraph import Paragraph
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        
        doc = Document(file_path)
        markdown_content = []
        
        # Text and tables
        for element in doc.element.body:
            if isinstance(element, CT_P):
                paragraph = Paragraph(element, doc)
                text = paragraph.text.strip()
                
                if text:
                    # Check for heading styles
                    if paragraph.style.name.startswith('Heading'):
                        level = paragraph.style.name.replace('Heading ', '')
                        try:
                            level_num = int(level)
                            markdown_content.append('#' * level_num + ' ' + text)
                        except ValueError:
                            markdown_content.append('## ' + text)
                    # Check for list items
                    elif paragraph.style.name.startswith('List'):
                        markdown_content.append('- ' + text)
                    else:
                        markdown_content.append(text)
                else:
                    markdown_content.append('')
                    
            elif isinstance(element, CT_Tbl):
                table = Table(element, doc)
                markdown_content.append('\n')
                # Convert table to markdown
                for i, row in enumerate(table.rows):
                    cells = [cell.text.strip() for cell in row.cells]
                    markdown_content.append('| ' + ' | '.join(cells) + ' |')
                    if i == 0:  # Add separator after header
                        markdown_content.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
                markdown_content.append('\n')

        # Images (export and append as a separate section)
        if image_dir is not None:
            image_dir_path = Path(image_dir)
            image_dir_path.mkdir(parents=True, exist_ok=True)
            images_added = 0

            for rel in doc.part.rels.values():
                if rel.reltype == RT.IMAGE:
                    image_part = rel.target_part
                    ext = Path(image_part.partname).suffix.lstrip('.') or 'png'
                    images_added += 1
                    filename = f"image_{images_added}.{ext}"
                    image_path = image_dir_path / filename

                    with open(image_path, "wb") as f:
                        f.write(image_part.blob)

                    rel_path = f"{image_dir_path.name}/{filename}"
                    if images_added == 1:
                        markdown_content.append('\n## Images\n')
                    markdown_content.append(f"![Image {images_added}]({rel_path})")
        
        return '\n'.join(markdown_content)
    
    except ImportError:
        raise ImportError("python-docx is required for .docx files. Install it with: pip install python-docx")
    except Exception as e:
        raise Exception(f"Error converting DOCX file: {str(e)}")


def convert_pdf_to_markdown(file_path, image_dir=None):
    """Convert a PDF file (.pdf) to Markdown.

    If image_dir is provided, embedded images are exported there and referenced
    as Markdown images grouped by page.
    """
    try:
        import fitz  # PyMuPDF
        
        doc = fitz.open(file_path)
        markdown_content = []
        image_dir_path = Path(image_dir) if image_dir is not None else None
        image_counter = 0
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            
            if text.strip():
                # Add page separator (except for first page)
                if page_num > 0:
                    markdown_content.append('\n---\n')
                
                # Process text and preserve basic formatting
                lines = text.split('\n')
                for line in lines:
                    line = line.strip()
                    if line:
                        # Try to detect headings (all caps or short lines)
                        if line.isupper() and len(line) < 100:
                            markdown_content.append('## ' + line)
                        else:
                            markdown_content.append(line)
                    else:
                        markdown_content.append('')

            # Export images on this page (if requested)
            if image_dir_path is not None:
                images = page.get_images(full=True)
                if images:
                    image_dir_path.mkdir(parents=True, exist_ok=True)
                    for img_index, img in enumerate(images, 1):
                        xref = img[0]
                        img_dict = doc.extract_image(xref)
                        ext = img_dict.get("ext", "png")
                        image_counter += 1
                        filename = f"page{page_num + 1}_image{img_index}.{ext}"
                        image_path = image_dir_path / filename
                        with open(image_path, "wb") as f:
                            f.write(img_dict["image"])
                        rel_path = f"{image_dir_path.name}/{filename}"
                        markdown_content.append(f"![Page {page_num + 1} Image {img_index}]({rel_path})")
        
        doc.close()
        return '\n'.join(markdown_content)
    
    except ImportError:
        raise ImportError("PyMuPDF is required for .pdf files. Install it with: pip install PyMuPDF")
    except Exception as e:
        raise Exception(f"Error converting PDF file: {str(e)}")


def convert_pptx_to_markdown(file_path, image_dir=None):
    """Convert a PowerPoint presentation (.pptx) to Markdown.

    If image_dir is provided, slide images are exported there and referenced
    as Markdown images near their slide content.
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        markdown_content = []
        image_dir_path = Path(image_dir) if image_dir is not None else None
        image_counter = 0
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Add slide header
            markdown_content.append(f'\n## Slide {slide_num}\n')
            
            # Process all shapes in the slide
            for shape in slide.shapes:
                # Tables -> Markdown tables
                if hasattr(shape, "has_table") and shape.has_table:
                    table = shape.table
                    markdown_content.append('')
                    for row_idx, row in enumerate(table.rows):
                        cells = [cell.text.replace('\n', ' ').strip() for cell in row.cells]
                        markdown_content.append('| ' + ' | '.join(cells) + ' |')
                        if row_idx == 0:
                            markdown_content.append('| ' + ' | '.join(['---'] * len(cells)) + ' |')
                    markdown_content.append('')
                    # Skip generic text handling for table shapes
                    continue

                # Charts -> Markdown tables of chart data (best effort)
                if hasattr(shape, "has_chart") and shape.has_chart:
                    chart = shape.chart
                    markdown_content.append('')

                    # Get categories if available
                    categories = []
                    try:
                        plot = chart.plots[0]
                        if plot.categories is not None:
                            for c in plot.categories:
                                # category objects may have .label or be basic types
                                label = getattr(c, "label", None)
                                categories.append(str(label if label is not None else c))
                    except Exception:
                        categories = []

                    # Build header: Category + each series name
                    headers = ['Category']
                    for series in chart.series:
                        headers.append(series.name if series.name is not None else 'Series')

                    markdown_content.append('| ' + ' | '.join(headers) + ' |')
                    markdown_content.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')

                    # Build rows
                    num_points = 0
                    if chart.series:
                        num_points = len(chart.series[0].values)

                    for idx in range(num_points):
                        if categories and idx < len(categories):
                            row = [categories[idx]]
                        else:
                            row = [f'Point {idx + 1}']

                        for series in chart.series:
                            try:
                                val = series.values[idx]
                            except Exception:
                                val = ''
                            row.append(str(val))

                        markdown_content.append('| ' + ' | '.join(row) + ' |')

                    markdown_content.append('')
                    # Skip generic text handling for chart shapes
                    continue

                # Export picture shapes as images
                if image_dir_path is not None and hasattr(shape, "image"):
                    image = shape.image
                    if image is not None:
                        image_dir_path.mkdir(parents=True, exist_ok=True)
                        image_counter += 1
                        ext = image.ext or "png"
                        filename = f"slide{slide_num}_image{image_counter}.{ext}"
                        image_path = image_dir_path / filename
                        with open(image_path, "wb") as f:
                            f.write(image.blob)
                        rel_path = f"{image_dir_path.name}/{filename}"
                        markdown_content.append(f"![Slide {slide_num} Image {image_counter}]({rel_path})")

                # Text handling
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Check if it's a title (usually first shape or specific shape type)
                    if shape == slide.shapes[0] or (hasattr(shape, 'is_placeholder') and 
                                                     shape.is_placeholder and 
                                                     shape.placeholder_format.idx == 0):
                        markdown_content.append('### ' + text)
                    else:
                        # Check for bullet points
                        if hasattr(shape, 'text_frame'):
                            for paragraph in shape.text_frame.paragraphs:
                                para_text = paragraph.text.strip()
                                if para_text:
                                    # Check if it's a bullet point
                                    if paragraph.level > 0 or para_text.startswith('•') or para_text.startswith('-'):
                                        indent = '  ' * paragraph.level
                                        markdown_content.append(indent + '- ' + para_text.lstrip('•- '))
                                    else:
                                        markdown_content.append(para_text)
                        else:
                            markdown_content.append(text)
            
            markdown_content.append('')  # Add blank line between slides
        
        return '\n'.join(markdown_content)
    
    except ImportError:
        raise ImportError("python-pptx is required for .pptx files. Install it with: pip install python-pptx")
    except Exception as e:
        raise Exception(f"Error converting PPTX file: {str(e)}")


def convert_ppt_to_markdown(file_path, image_dir=None):
    """Convert a legacy PowerPoint presentation (.ppt) to Markdown.

    This uses Microsoft PowerPoint via COM automation to first convert the
    .ppt file to .pptx, then reuses the existing .pptx conversion logic.
    """
    try:
        import win32com.client
        import pythoncom

        # Initialize COM
        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")

        temp_pptx_path = None
        abs_path = str(Path(file_path).resolve())
        try:
            # Use absolute path and open with a window to avoid visibility errors
            presentation = ppt_app.Presentations.Open(abs_path, WithWindow=True)
            try:
                # Create a temporary .pptx file path
                temp_dir = tempfile.gettempdir()
                temp_pptx_path = os.path.join(
                    temp_dir, f"ppt_convert_{uuid.uuid4().hex}.pptx"
                )

                # 24 = ppSaveAsOpenXMLPresentation
                presentation.SaveAs(temp_pptx_path, 24)
            finally:
                presentation.Close()

            # Reuse the existing .pptx converter
            markdown_content = convert_pptx_to_markdown(temp_pptx_path, image_dir=image_dir)
        finally:
            ppt_app.Quit()
            pythoncom.CoUninitialize()
            if temp_pptx_path and os.path.exists(temp_pptx_path):
                try:
                    os.remove(temp_pptx_path)
                except OSError:
                    # If cleanup fails, it's not critical for the conversion result
                    pass

        return markdown_content

    except ImportError:
        raise ImportError(
            "pywin32 is required for .ppt files, or convert the file to .pptx manually. "
            "Install it with: pip install pywin32"
        )
    except Exception as e:
        raise Exception(
            "Error converting PPT file. Make sure Microsoft PowerPoint is installed "
            f"and accessible. Details: {str(e)}"
        )


def convert_file_to_markdown(input_file, output_file=None):
    """
    Convert a file to Markdown format.
    
    Args:
        input_file: Path to the input file
        output_file: Optional path to the output file. If not provided, 
                     output will be saved as input_file.md
    
    Returns:
        Path to the output markdown file
    """
    input_path = Path(input_file)
    
    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_file}")
    
    # Determine output file path
    if output_file is None:
        output_path = input_path.with_suffix('.md')
    else:
        output_path = Path(output_file)

    # Directory where extracted images will be stored (next to the markdown file)
    image_dir = output_path.parent / f"{output_path.stem}_images"
    
    # Get file extension
    ext = input_path.suffix.lower()
    
    # Convert based on file type
    print(f"Converting {input_path.name} to Markdown...")
    
    if ext == '.docx':
        markdown_content = convert_docx_to_markdown(input_path, image_dir=image_dir)
    elif ext == '.pdf':
        markdown_content = convert_pdf_to_markdown(input_path, image_dir=image_dir)
    elif ext == '.pptx':
        markdown_content = convert_pptx_to_markdown(input_path, image_dir=image_dir)
    elif ext == '.ppt':
        markdown_content = convert_ppt_to_markdown(input_path, image_dir=image_dir)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Supported formats: .docx, .pdf, .ppt, .pptx")
    
    # Write markdown file
    with open(output_path, 'w', encoding='utf-8') as f:
        f.write(markdown_content)
    
    print(f"Successfully converted to: {output_path}")
    return str(output_path)


def main():
    """Main entry point for the command-line interface."""
    parser = argparse.ArgumentParser(
        description='Convert Word, PDF, and PowerPoint files to Markdown format',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python file_to_markdown.py document.docx
  python file_to_markdown.py presentation.pdf -o output.md
  python file_to_markdown.py slides.pptx
        """
    )
    
    parser.add_argument('input_file', help='Path to the input file (.docx, .pdf, or .pptx)')
    parser.add_argument('-o', '--output', dest='output_file', 
                       help='Path to the output markdown file (default: input_file.md)')
    
    args = parser.parse_args()
    
    try:
        convert_file_to_markdown(args.input_file, args.output_file)
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == '__main__':
    main()
